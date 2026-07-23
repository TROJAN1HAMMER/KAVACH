"""
KAVACH pitch deck -- shared visual system for python-pptx generation.

Dark banking-blue theme, consulting-deck layout primitives (kicker/title blocks,
glass cards, process-flow chevrons, comparison tables, screenshot placeholders).
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

# ---------------------------------------------------------------- geometry --
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.62)
CONTENT_BOTTOM = Inches(6.95)
FOOTER_Y = Inches(7.08)

# ----------------------------------------------------------------- palette --
BG        = RGBColor(0x0A, 0x0E, 0x1A)
BG_SOFT   = RGBColor(0x0D, 0x12, 0x20)
PANEL     = RGBColor(0x13, 0x1B, 0x2E)
PANEL_ALT = RGBColor(0x18, 0x23, 0x3B)
PANEL_HI  = RGBColor(0x1E, 0x2B, 0x47)
BORDER    = RGBColor(0x2A, 0x36, 0x54)
BORDER_LT = RGBColor(0x3A, 0x48, 0x6B)

GOLD      = RGBColor(0xED, 0xA1, 0x00)
GOLD_LT   = RGBColor(0xFF, 0xC6, 0x4B)
BLUE      = RGBColor(0x2A, 0x78, 0xD6)
BLUE_LT   = RGBColor(0x6C, 0xAE, 0xF2)
CYAN      = RGBColor(0x38, 0xBD, 0xF8)
RED       = RGBColor(0xE3, 0x49, 0x48)
RED_LT    = RGBColor(0xF2, 0x8B, 0x8A)
PURPLE    = RGBColor(0x4A, 0x3A, 0xA7)
PURPLE_LT = RGBColor(0x9B, 0x8B, 0xEA)
GREEN     = RGBColor(0x1B, 0xAF, 0x7A)
GREEN_LT  = RGBColor(0x6F, 0xDD, 0xB4)

WHITE     = RGBColor(0xF5, 0xF7, 0xFA)
SLATE     = RGBColor(0x9A, 0xA7, 0xBD)
SLATE_DIM = RGBColor(0x64, 0x70, 0x8A)

FONT = "Segoe UI"
FONT_LIGHT = "Segoe UI Light"
FONT_SEMI = "Segoe UI Semibold"

SEVERITY_COLOR = {
    "CRITICAL": RED,
    "HIGH": RGBColor(0xE8, 0x7A, 0x33),
    "MEDIUM": GOLD,
    "LOW": BLUE_LT,
}


# =============================================================== core i/o ==
def new_deck():
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _send_to_back(slide, shape):
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def add_slide(prs, section_label="", index=None, total=22):
    """Blank dark-navy slide with brand accent + footer chrome."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG, shape_type=MSO_SHAPE.RECTANGLE)
    _send_to_back(slide, bg)
    _brand_accent(slide)
    if index is not None:
        add_footer(slide, index, total, section_label)
    return slide


def _brand_accent(slide):
    """Thin gold arc motif bleeding off the right edge -- echoes the KAVACH brand mark."""
    d = Inches(6.4)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, SLIDE_W - Inches(2.0), Inches(-3.6), d, d)
    circ.fill.background()
    circ.line.color.rgb = GOLD
    circ.line.width = Pt(1.1)
    set_line_alpha(circ, 34)
    no_shadow(circ)


def add_footer(slide, index, total, section_label=""):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, MARGIN, FOOTER_Y, SLIDE_W - MARGIN, FOOTER_Y)
    line.line.color.rgb = BORDER
    line.line.width = Pt(0.75)
    no_shadow(line)
    textbox(slide, MARGIN, Inches(7.14), Inches(5.0), Inches(0.3),
            "KAVACH  ·  PSB Hackathon 2026", size=9, color=SLATE_DIM, font=FONT)
    if section_label:
        textbox(slide, Inches(4.9), Inches(7.14), Inches(5.0), Inches(0.3), section_label.upper(),
                size=9, color=SLATE_DIM, align=PP_ALIGN.CENTER, font=FONT)
    textbox(slide, SLIDE_W - Inches(1.5), Inches(7.14), Inches(1.0), Inches(0.3),
            f"{index:02d} / {total}", size=9, color=SLATE_DIM, align=PP_ALIGN.RIGHT, font=FONT)


# =============================================================== shapes ===
def no_shadow(shape):
    shape.shadow.inherit = False


def set_alpha(shape, pct):
    """Opacity 0-100 (100 = opaque) on a shape's solid fore-color fill."""
    try:
        solidFill = shape.fill.fore_color._xFill
    except Exception:
        return
    clr = solidFill.find(qn('a:srgbClr'))
    if clr is None:
        return
    for a in clr.findall(qn('a:alpha')):
        clr.remove(a)
    alpha = OxmlElement('a:alpha')
    alpha.set('val', str(int(pct * 1000)))
    clr.append(alpha)


def set_line_alpha(shape, pct):
    try:
        solidFill = shape.line.color._xFill
    except Exception:
        return
    clr = solidFill.find(qn('a:srgbClr'))
    if clr is None:
        return
    for a in clr.findall(qn('a:alpha')):
        clr.remove(a)
    alpha = OxmlElement('a:alpha')
    alpha.set('val', str(int(pct * 1000)))
    clr.append(alpha)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75,
         shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, shadow=False, dashed=False):
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
        if dashed:
            shp.line.dash_style = None
            ln = shp.line._get_or_add_ln()
            pd = OxmlElement('a:prstDash')
            pd.set('val', 'dash')
            ln.append(pd)
    if not shadow:
        no_shadow(shp)
    return shp


def card(slide, x, y, w, h, fill=PANEL, line=BORDER, radius=0.055, alpha=None, line_w=1.0):
    shp = rect(slide, x, y, w, h, fill=fill, line=line, radius=radius, line_w=line_w)
    if alpha is not None:
        set_alpha(shp, alpha)
    return shp


def hline(slide, x, y, w, color=BORDER, weight=1.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    no_shadow(ln)
    return ln


def vline(slide, x, y, h, color=BORDER, weight=1.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x, y + h)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    no_shadow(ln)
    return ln


def pill(slide, x, y, w, h, text, fill=GOLD, text_color=BG, size=10.5, bold=True, line=None):
    shp = rect(slide, x, y, w, h, fill=fill, line=line, radius=0.5, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = FONT
    return shp


def dot(slide, cx, cy, d, fill):
    shp = rect(slide, cx - d / 2, cy - d / 2, d, d, fill=fill, shape_type=MSO_SHAPE.OVAL, line=None)
    return shp


# =============================================================== text ====
def textbox(slide, x, y, w, h, text, size=14, color=WHITE, bold=False, italic=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.0,
            wrap=True, space_after=0, letter_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
        if letter_spacing is not None:
            _set_letter_spacing(run, letter_spacing)
    return tb


def _set_letter_spacing(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(pts * 100)))


def rich_line(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    """Single paragraph, multiple differently-styled runs. runs = list of dicts."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for r in runs:
        run = p.add_run()
        run.text = r.get('text', '')
        run.font.size = Pt(r.get('size', 14))
        run.font.bold = r.get('bold', False)
        run.font.italic = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
        run.font.name = r.get('font', FONT)
    return tb


def kicker_title(slide, kicker, title, subtitle=None, title_size=30, y=Inches(0.5), title_box_w=11.8):
    """Eyebrow + big title + thin gold rule + optional subtitle. Standard slide header.
    Estimates wrapped line count so the divider never collides with a 2-line title."""
    textbox(slide, MARGIN, y, Inches(10.5), Inches(0.32), kicker.upper(),
            size=12.5, color=GOLD, bold=True, font=FONT, letter_spacing=1.6)
    title_y = y + Inches(0.34)
    avg_char_w_in = title_size * 0.54 / 72.0
    chars_per_line = max(10, int(title_box_w / avg_char_w_in))
    n_lines = max(1, -(-len(title) // chars_per_line))  # ceil div
    line_h_in = title_size * 1.22 / 72.0
    title_h_in = max(0.62, n_lines * line_h_in + 0.06)
    textbox(slide, MARGIN, title_y, Inches(title_box_w), Inches(title_h_in), title,
            size=title_size, color=WHITE, bold=True, font=FONT, line_spacing=1.02)
    ty = title_y + Inches(title_h_in)
    if subtitle:
        textbox(slide, MARGIN, ty, Inches(title_box_w), Inches(0.36), subtitle,
                size=13.5, color=SLATE, italic=True, font=FONT)
        ty += Inches(0.36)
    hline(slide, MARGIN, ty + Inches(0.07), Inches(1.05), color=GOLD, weight=2.25)
    return ty + Inches(0.07)


def bullets(slide, x, y, w, h, items, size=13, color=WHITE, marker_color=GOLD,
            gap=0.12, marker="chevron", line_spacing=1.05, bold_lead=False):
    """items: list of str OR (lead, rest) tuples for a bold-lead bullet."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(gap * 72)
        m = p.add_run()
        m.text = ("›  " if marker == "chevron" else "•  ")
        m.font.size = Pt(size)
        m.font.bold = True
        m.font.color.rgb = marker_color
        m.font.name = FONT
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = lead
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = WHITE
            r1.font.name = FONT
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = FONT
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = FONT
    return tb


# ========================================================= composed cards =
def stat_tile(slide, x, y, w, h, value, label, accent=GOLD, sub=None):
    c = card(slide, x, y, w, h, fill=PANEL, line=BORDER)
    hline(slide, x + Inches(0.18), y + Inches(0.16), Inches(0.34), color=accent, weight=2.5)
    textbox(slide, x + Inches(0.18), y + Inches(0.26), w - Inches(0.36), Inches(0.56),
            value, size=27, color=WHITE, bold=True, font=FONT)
    sub_h = Inches(0.4) if sub else Inches(0.0)
    label_top = y + Inches(0.9)
    label_h = h - Inches(0.9) - Inches(0.14) - sub_h
    textbox(slide, x + Inches(0.18), label_top, w - Inches(0.36), label_h,
            label, size=10.5, color=SLATE, font=FONT, line_spacing=1.18)
    if sub:
        textbox(slide, x + Inches(0.18), y + h - sub_h - Inches(0.12), w - Inches(0.36), sub_h,
                sub, size=8, color=SLATE_DIM, italic=True, font=FONT, anchor=MSO_ANCHOR.BOTTOM)
    return c


def mini_stat(slide, x, y, w, h, value, label, accent=GOLD):
    """Compact single-row stat banner for short tiles (h ~0.7-1.0in)."""
    card(slide, x, y, w, h, fill=PANEL, line=BORDER)
    vline(slide, x + Inches(0.02), y + Inches(0.14), h - Inches(0.28), color=accent, weight=3)
    textbox(slide, x + Inches(0.24), y + Inches(0.08), w - Inches(0.44), Inches(0.32),
            value, size=15, color=accent, bold=True)
    textbox(slide, x + Inches(0.24), y + Inches(0.4), w - Inches(0.44), h - Inches(0.48),
            label, size=9, color=SLATE, line_spacing=1.05)


def icon_card(slide, x, y, w, h, title, body, accent=GOLD, title_size=13.5, body_size=11):
    card(slide, x, y, w, h, fill=PANEL, line=BORDER)
    dot(slide, x + Inches(0.28), y + Inches(0.30), Inches(0.14), accent)
    textbox(slide, x + Inches(0.5), y + Inches(0.19), w - Inches(0.7), Inches(0.32),
            title, size=title_size, color=WHITE, bold=True, font=FONT)
    textbox(slide, x + Inches(0.24), y + Inches(0.58), w - Inches(0.48), h - Inches(0.75),
            body, size=body_size, color=SLATE, font=FONT, line_spacing=1.08)


def process_flow(slide, steps, x, y, w, h, colors=None, label_size=11.5, num=True, gap_in=0.16):
    """Horizontal chevron pipeline. steps: list[str]."""
    n = len(steps)
    gap = Inches(gap_in)
    box_w = Emu(int((w - gap * (n - 1)) / n))
    colors = colors or [GOLD] * n
    for i, s in enumerate(steps):
        bx = x + i * (box_w + gap)
        shape_type = MSO_SHAPE.CHEVRON if 0 < i < n else MSO_SHAPE.PENTAGON
        c = rect(slide, bx, y, box_w, h, fill=PANEL_ALT, line=colors[i % len(colors)],
                 line_w=1.25, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(6); tf.margin_right = Pt(6)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = s
        r.font.size = Pt(label_size)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = FONT
        if i < n - 1:
            ax = bx + box_w + Inches(0.015)
            arrow_w = gap - Inches(0.03)
            tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, ax, y + h / 2 - Inches(0.07),
                                          arrow_w, Inches(0.14))
            tri.rotation = 90
            tri.fill.solid(); tri.fill.fore_color.rgb = GOLD
            tri.line.fill.background()
            no_shadow(tri)


def numbered_node(slide, x, y, w, h, number, title, body, accent=GOLD, body_size=10.5):
    card(slide, x, y, w, h, fill=PANEL, line=BORDER)
    circ = rect(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.42), Inches(0.42),
                fill=accent, line=None, shape_type=MSO_SHAPE.OVAL)
    tf = circ.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(number); r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = BG; r.font.name = FONT
    textbox(slide, x + Inches(0.76), y + Inches(0.24), w - Inches(0.95), Inches(0.4),
            title, size=13, color=WHITE, bold=True, font=FONT)
    textbox(slide, x + Inches(0.22), y + Inches(0.75), w - Inches(0.42), h - Inches(0.9),
            body, size=body_size, color=SLATE, font=FONT, line_spacing=1.08)


def arch_node(slide, x, y, w, h, label, sub=None, accent=GOLD, size=10.5):
    c = card(slide, x, y, w, h, fill=PANEL_ALT, line=accent, line_w=1.25, radius=0.16)
    tf = c.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(5); tf.margin_right = Pt(5); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = FONT
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(size - 2.5); r2.font.color.rgb = SLATE
        r2.font.name = FONT
    return c


def connector(slide, x1, y1, x2, y2, color=BORDER_LT, weight=1.25, dashed=False, arrow=True):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    no_shadow(ln)
    if dashed:
        lnEl = ln.line._get_or_add_ln()
        pd = OxmlElement('a:prstDash'); pd.set('val', 'dash')
        lnEl.append(pd)
    if arrow:
        lnEl = ln.line._get_or_add_ln()
        tail = OxmlElement('a:tailEnd')
        tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
        lnEl.append(tail)
    return ln


def screenshot_placeholder(slide, x, y, w, h, label, route=None, note=None):
    """Browser-chrome style dashed placeholder for a not-yet-captured UI screenshot."""
    frame = card(slide, x, y, w, h, fill=PANEL, line=BORDER_LT, radius=0.035)
    bar = rect(slide, x, y, w, Inches(0.3), fill=PANEL_HI, line=None, radius=0.0,
               shape_type=MSO_SHAPE.RECTANGLE)
    for i in range(3):
        dot(slide, x + Inches(0.18 + i * 0.16), y + Inches(0.15), Inches(0.09),
            [RED, GOLD, GREEN][i])
    if route:
        textbox(slide, x + Inches(0.6), y + Inches(0.05), w - Inches(0.8), Inches(0.22),
                route, size=8.5, color=SLATE_DIM, font="Consolas")
    inner = rect(slide, x + Inches(0.16), y + Inches(0.44), w - Inches(0.32), h - Inches(0.6),
                 fill=None, line=SLATE_DIM, radius=0.04, dashed=True, line_w=1.0)
    tf = inner.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "▣  SCREENSHOT PLACEHOLDER"
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = SLATE
    r.font.name = FONT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(10.5); r2.font.color.rgb = GOLD_LT; r2.font.name = FONT
    if note:
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(4)
        r3 = p3.add_run(); r3.text = note
        r3.font.size = Pt(8.5); r3.font.italic = True; r3.font.color.rgb = SLATE_DIM
        r3.font.name = FONT
    return frame


def shield_mark(slide, x, y, size, fill=GOLD, inner=BG):
    """Simple two-tone shield glyph built from a freeform outline."""
    w = size
    h = size * Inches(1) / Inches(1)  # keep 1:1 ratio scaling handled by caller
    fb = slide.shapes.build_freeform(start_x=x + w * 0.5, start_y=y, scale=1)
    pts = [
        (x + w, y + h * 0.12),
        (x + w, y + h * 0.5),
        (x + w * 0.5, y + h),
        (x, y + h * 0.5),
        (x, y + h * 0.12),
        (x + w * 0.5, y),
    ]
    fb2 = slide.shapes.build_freeform(start_x=int(x + w * 0.5), start_y=int(y), scale=1)
    fb2.add_line_segments([(int(px), int(py)) for px, py in pts], close=True)
    shp = fb2.convert_to_shape()
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    no_shadow(shp)
    return shp


def source_note(slide, x, y, w, text, size=8.5):
    textbox(slide, x, y, w, Inches(0.3), text, size=size, color=SLATE_DIM, italic=True, font=FONT)


def set_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


# ==================================================================table ==
def themed_table(slide, x, y, w, h, headers, rows, col_widths=None,
                  header_fill=PANEL_HI, header_color=GOLD, body_size=11.5,
                  header_size=11.5, row_fill_a=PANEL, row_fill_b=BG_SOFT,
                  highlight_col=None, highlight_color=GOLD, first_col_bold=True):
    """Native pptx table themed for the dark deck. rows: list[list[str]]."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gframe = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = gframe.table
    # kill the built-in banded style so our fills show cleanly
    tbl = table._tbl
    from pptx.oxml.ns import qn as _qn
    tblPr = tbl.find(_qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')
        for child in list(tblPr):
            tblPr.remove(child)

    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    table.rows[0].height = Emu(int(h / n_rows * 0.9))
    for i in range(1, n_rows):
        table.rows[i].height = Emu(int(h / n_rows))

    for j, htext in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.margin_left = Pt(8); cell.margin_right = Pt(8)
        cell.margin_top = Pt(4); cell.margin_bottom = Pt(4)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = htext
        r.font.size = Pt(header_size); r.font.bold = True
        r.font.color.rgb = header_color if j != (highlight_col or -99) else header_color
        r.font.name = FONT

    for i, row in enumerate(rows):
        fill = row_fill_a if i % 2 == 0 else row_fill_b
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            is_hl = (highlight_col is not None and j == highlight_col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill if not is_hl else PANEL_HI
            cell.margin_left = Pt(8); cell.margin_right = Pt(8)
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            r.font.size = Pt(body_size)
            r.font.name = FONT
            r.font.color.rgb = (GOLD_LT if is_hl else (WHITE if (j == 0 and first_col_bold) else SLATE))
            r.font.bold = is_hl or (j == 0 and first_col_bold)
    return gframe
