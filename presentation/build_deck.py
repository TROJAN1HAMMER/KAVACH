"""
Builds KAVACH_PSB_Hackathon_2026.pptx -- a 22-slide investor/jury deck grounded
strictly in the actual KAVACH repository implementation (backend/, frontend/,
mobile/, docs/). Run with the Python 3.12 interpreter that has python-pptx
installed (see README at top of repo for the exact path used).
"""
import os
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from deck_helpers import (
    new_deck, add_slide, card, rect, textbox, rich_line, kicker_title, bullets,
    stat_tile, mini_stat, icon_card, process_flow, numbered_node, arch_node, connector,
    screenshot_placeholder, shield_mark, source_note, set_notes, themed_table,
    pill, hline, vline, dot, no_shadow,
    BG, BG_SOFT, PANEL, PANEL_ALT, PANEL_HI, BORDER, BORDER_LT,
    GOLD, GOLD_LT, BLUE, BLUE_LT, CYAN, RED, RED_LT, PURPLE, PURPLE_LT,
    GREEN, GREEN_LT, WHITE, SLATE, SLATE_DIM, FONT, SLIDE_W, SLIDE_H, MARGIN,
)

TOTAL = 22
TEAM_NAME = "Kavach"
MEMBERS = ["Harshith B", "Tejashwan Gangishetty", "Pratham Lal"]
INSTITUTE = "Vellore Institute of Technology"
GITHUB = "github.com/TROJAN1HAMMER/KAVACH"

CW = 12.23          # content width (margin to margin), inches
LEFT = 0.55


# ------------------------------------------------------------------ 01 ----
def slide_01_title(prs):
    s = add_slide(prs, "", 1, TOTAL)
    shield_mark(s, Inches(0.9), Inches(0.72), Inches(0.72), fill=GOLD, inner=BG)
    rich_line(s, Inches(1.78), Inches(0.66), Inches(6), Inches(0.7),
              [{"text": "KAVACH", "size": 40, "bold": True, "color": WHITE, "font": FONT}])
    textbox(s, Inches(1.8), Inches(1.32), Inches(6), Inches(0.32),
            "S E C U R E   ·   A N A L Y Z E   ·   P R O T E C T", size=11.5, color=GOLD, bold=True)
    hline(s, Inches(0.9), Inches(1.86), Inches(9.2), color=BORDER_LT, weight=1)

    textbox(s, Inches(0.9), Inches(2.35), Inches(11.2), Inches(0.9),
            "AI-Powered DevSecOps Platform for Banking", size=33, color=WHITE, bold=True)
    textbox(s, Inches(0.92), Inches(3.15), Inches(10.6), Inches(0.5),
            "Static code analysis, dependency & configuration scanning, business-aware risk scoring, "
            "regulatory compliance mapping, and grounded AI -- across one deduplicated pipeline.",
            size=14, color=SLATE, italic=True)

    pill(s, Inches(0.92), Inches(3.85), Inches(2.05), Inches(0.42), "PSB HACKATHON 2026", fill=GOLD, text_color=BG, size=12)
    textbox(s, Inches(3.1), Inches(3.85), Inches(9.1), Inches(0.42),
            "Problem Statement 3 -- Automated Discovery of Misconfigurations in Open-Source Dependencies",
            size=12.5, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)

    # team block
    tb = card(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(1.55), fill=PANEL, line=BORDER)
    colw = 11.5 / 3
    labels = [("TEAM", TEAM_NAME), ("TEAM MEMBERS", "\n".join(MEMBERS)), ("INSTITUTE", INSTITUTE)]
    for i, (lab, val) in enumerate(labels):
        x = 0.9 + i * colw
        if i > 0:
            vline(s, Inches(x), Inches(4.75), Inches(1.15), color=BORDER_LT)
        textbox(s, Inches(x + 0.3), Inches(4.75), Inches(colw - 0.5), Inches(0.28), lab,
                size=10.5, color=GOLD, bold=True, letter_spacing=1.2)
        textbox(s, Inches(x + 0.3), Inches(5.08), Inches(colw - 0.5), Inches(0.95), val,
                size=13.5 if i != 1 else 12.5, color=WHITE, bold=(i != 1), line_spacing=1.25)

    textbox(s, Inches(0.9), Inches(6.35), Inches(9), Inches(0.35), GITHUB, size=11, color=SLATE_DIM)
    set_notes(s, "Welcome the jury. KAVACH is a working, multi-service DevSecOps platform we built for PSB "
                 "Hackathon 2026's Problem Statement 3 -- automated discovery of misconfigurations in "
                 "open-source dependencies for banking applications. Everything in this deck maps to code "
                 "that exists in the repository today; where something is a placeholder or future work, we "
                 "say so explicitly.")


# ------------------------------------------------------------------ 02 ----
def slide_02_about(prs):
    s = add_slide(prs, "About the Hackathon", 2, TOTAL)
    kicker_title(s, "PSB Hackathon 2026", "A National Mission for Bank-Grade Cybersecurity")

    y = 2.05
    # hierarchy diagram
    arch_node(s, Inches(LEFT), Inches(y), Inches(3.6), Inches(0.62),
              "Department of Financial Services", "Ministry of Finance, Government of India", accent=BLUE, size=11)
    arch_node(s, Inches(LEFT + 3.85), Inches(y), Inches(3.6), Inches(0.62),
              "Indian Banks' Association", "Sector-wide banking industry body", accent=BLUE, size=11)
    connector(s, Inches(LEFT + 1.8), Inches(y + 0.62), Inches(LEFT + 3.9), Inches(y + 1.05), color=BORDER_LT)
    connector(s, Inches(LEFT + 5.65), Inches(y + 0.62), Inches(LEFT + 3.9 + 1.9), Inches(y + 1.05), color=BORDER_LT)
    arch_node(s, Inches(LEFT + 2.0), Inches(y + 1.05), Inches(3.9), Inches(0.62), "PSB Hackathon 2026",
              "National-level initiative", accent=GOLD, size=12)
    connector(s, Inches(LEFT + 3.95), Inches(y + 1.67), Inches(LEFT + 1.6), Inches(y + 2.1), color=BORDER_LT)
    connector(s, Inches(LEFT + 3.95), Inches(y + 1.67), Inches(LEFT + 6.3), Inches(y + 2.1), color=BORDER_LT)
    arch_node(s, Inches(LEFT), Inches(y + 2.1), Inches(3.6), Inches(0.62), "Hosted by UCO Bank",
              "Public Sector Bank -- problem owner", accent=GREEN, size=11)
    arch_node(s, Inches(LEFT + 3.85), Inches(y + 2.1), Inches(3.6), Inches(0.62), "Hosted by IIT Kharagpur",
              "Academic & technical partner", accent=GREEN, size=11)

    icon_card(s, Inches(8.1), Inches(2.05), Inches(4.13), Inches(1.55), "Mission",
              "Strengthening cybersecurity across India's Public Sector Banks -- turning "
              "hackathon-grade prototypes into evaluable, real security tooling.", accent=GOLD)
    icon_card(s, Inches(8.1), Inches(3.75), Inches(4.13), Inches(1.98), "Why This Challenge Matters",
              "PSBs run large hybrid stacks (legacy core banking + modern microservices) built on hundreds "
              "of open-source components, under RBI's IT Framework (2021), PCI DSS, and SWIFT CSP "
              "obligations -- with security teams too thin to manually triage every release.",
              accent=BLUE)

    set_notes(s, "PSB Hackathon 2026 is led by the Department of Financial Services under the Ministry of "
                 "Finance, together with the Indian Banks' Association, and hosted by UCO Bank with IIT "
                 "Kharagpur as academic partner. The mission is explicit: strengthen cybersecurity across "
                 "Public Sector Banks. That framing matters for how we built KAVACH -- every design choice, "
                 "from the Banking Risk Score to RBI/PCI/SWIFT compliance mapping, is aimed at a PSB's real "
                 "operating constraints, not a generic enterprise.")


# ------------------------------------------------------------------ 03 ----
def slide_03_problem(prs):
    s = add_slide(prs, "Problem Statement", 3, TOTAL)
    kicker_title(s, "Problem Statement 3", "Automated Discovery of Misconfigurations in Open-Source Dependencies")

    card(s, Inches(LEFT), Inches(2.1), Inches(5.95), Inches(4.55), fill=PANEL, line=BORDER)
    textbox(s, Inches(LEFT + 0.28), Inches(2.28), Inches(5.4), Inches(0.32), "THE GAP TODAY",
            size=12, color=GOLD, bold=True, letter_spacing=1.2)
    bullets(s, Inches(LEFT + 0.28), Inches(2.66), Inches(5.45), Inches(3.9), [
        ("Hundreds of OSS libraries ", "per banking application -- often unpatched or misconfigured."),
        ("Known-CVE tools only. ", "Existing SCA scanners catch published CVEs, not insecure defaults: "
                                     "hardcoded passwords, weak cipher suites, insecure configs, unsafe "
                                     "dependency usage."),
        ("Manual compliance mapping. ", "Turning a raw finding into “this violates PCI DSS 6.2” is "
                                          "normally a spreadsheet exercise, after the fact."),
        ("Release pressure wins. ", "Developers miss vulnerabilities under rapid release cycles -- the "
                                      "attack surface grows faster than it's reviewed."),
    ], size=13, gap=0.16)

    card(s, Inches(6.83), Inches(2.1), Inches(5.4), Inches(4.55), fill=PANEL, line=GOLD, line_w=1.25)
    textbox(s, Inches(7.11), Inches(2.28), Inches(4.9), Inches(0.32), "WHAT THE CHALLENGE ASKS FOR",
            size=12, color=GOLD, bold=True, letter_spacing=1.0)
    checks = ["Scan source code", "Scan dependencies", "Detect known CVEs",
              "Detect insecure configurations", "Integrate with CI/CD",
              "Produce clear, actionable reports", "Recommend automated remediation"]
    y = 2.7
    for c in checks:
        dot(s, Inches(7.32), Inches(y + 0.11), Inches(0.14), GREEN)
        textbox(s, Inches(7.55), Inches(y), Inches(4.5), Inches(0.34), c, size=13, color=WHITE)
        y += 0.475

    set_notes(s, "This is Problem Statement 3, verbatim in spirit: banking applications depend on hundreds "
                 "of open-source libraries; traditional SCA tools stop at known CVEs and miss insecure "
                 "defaults -- hardcoded passwords, weak cipher suites, insecure configuration, unsafe "
                 "dependency usage. The ask is a static analysis platform that scans code and dependencies, "
                 "detects both CVEs and misconfigurations, integrates with CI/CD, and produces clear, "
                 "actionable, remediation-oriented reports. The rest of this deck walks through exactly how "
                 "KAVACH answers each one of these, module by module, with evidence from the repository.")


# ------------------------------------------------------------------ 04 ----
def slide_04_industry(prs):
    s = add_slide(prs, "Industry Challenges", 4, TOTAL)
    kicker_title(s, "Why This Is Urgent", "The Threat Landscape, In Numbers")

    stats = [
        ("75%", "Growth in open-source malware discovered in 2025 vs. 2024 -- 1.23M+ malicious "
                 "packages now blocked across npm, PyPI, Maven, NuGet & Hugging Face.", GOLD,
         "Sonatype, 2026 State of the Software Supply Chain"),
        ("95%", "Of the time a vulnerable open-source component is pulled into a build, a fixed "
                 "version already existed -- this is a patching gap, not a detection gap.", BLUE_LT,
         "Sonatype, 10th Annual State of the Software Supply Chain (2024)"),
        ("37%", "Of breaches in 2025 involved exploited vulnerabilities, unpatched systems, or "
                 "misconfigurations.", RED_LT,
         "Verizon, 2025 Data Breach Investigations Report"),
        ("4.1M", "Average monthly cyberattacks against India's BFSI sector in 2025, with 248+ "
                  "confirmed breaches across scheduled commercial banks.", GREEN_LT,
         "CERT-In / CSIRT-Fin / SISA Digital Threat Report; 2025 industry reporting"),
    ]
    w = (CW - 0.18 * 3) / 4
    for i, (val, label, color, src) in enumerate(stats):
        x = LEFT + i * (w + 0.18)
        stat_tile(s, Inches(x), Inches(2.1), Inches(w), Inches(2.85), val, label, accent=color, sub=src)

    card(s, Inches(LEFT), Inches(5.15), Inches(CW), Inches(1.15), fill=PANEL_ALT, line=GOLD, line_w=1.1)
    textbox(s, Inches(LEFT + 0.3), Inches(5.35), Inches(CW - 0.6), Inches(0.75),
            "The bottleneck isn't detecting known CVEs -- it's catching what CVE databases don't cover: "
            "insecure defaults, weak configuration, and the lag between “a fix exists” and "
            "“a fix is deployed.”",
            size=14, color=WHITE, italic=True, anchor=MSO_ANCHOR.MIDDLE)

    set_notes(s, "Four sourced numbers, not vibes. Sonatype's 2026 supply-chain report shows open-source "
                 "malware up 75% year over year. Their 2024 report found that 95% of the time a vulnerable "
                 "component is consumed, a fix already existed -- so the real failure is patch lag, exactly "
                 "the insecure-defaults gap this problem statement calls out. Verizon's 2025 DBIR puts "
                 "exploited vulnerabilities, unpatched systems, and misconfigurations at 37% of breaches. "
                 "And closer to home, India's BFSI sector saw an average 4.1 million attacks a month in "
                 "2025 with 248+ confirmed bank breaches. Cite the sources on screen -- this is the kind of "
                 "slide a jury will fact-check.")


# ------------------------------------------------------------------ 05 ----
def slide_05_solution(prs):
    s = add_slide(prs, "Our Solution", 5, TOTAL)
    kicker_title(s, "Our Solution", "KAVACH: One Pipeline, From Repository to Boardroom")

    textbox(s, Inches(LEFT), Inches(2.05), Inches(7.1), Inches(1.0),
            "A distributed, AI-assisted DevSecOps platform that scans a repository across 9 independent "
            "security tools, deduplicates the results into one unified finding set, and turns that into "
            "what a bank can act on: a business-aware risk score, mapped compliance clauses, and a "
            "citation-backed explanation.",
            size=13.5, color=SLATE, line_spacing=1.2)

    process_flow(s, ["Repository", "CI / CD", "Scanning", "AI", "Risk Engine", "Compliance", "Reports"],
                 Inches(LEFT), Inches(3.35), Inches(7.1), Inches(0.85),
                 colors=[BLUE, BLUE, BLUE_LT, PURPLE, GOLD, GREEN, WHITE], label_size=8, gap_in=0.08)

    bullets(s, Inches(LEFT), Inches(4.55), Inches(7.1), Inches(2.2), [
        ("Single scan job, ", "9 scanners fanned out in parallel via a priority-queued Celery worker pool."),
        ("Cross-tool aggregation ", "collapses overlapping findings into one deduplicated, enriched "
                                     "UnifiedFinding set with CWE / OWASP / MITRE ATT&CK tags."),
        ("Business Risk Score + Attack Surface Exposure ", "quantify what to fix first and how exposed "
                                                              "the codebase is overall."),
        ("7 report formats ", "generated per scan: Executive PDF, Technical PDF, SARIF, CycloneDX SBOM, "
                                "unified JSON, compliance JSON, CSV."),
    ], size=12.5, gap=0.14)

    screenshot_placeholder(s, Inches(8.0), Inches(2.05), Inches(4.23), Inches(4.7),
                            "Dashboard Overview", route="/  (OverviewPage.tsx)",
                            note="Repo/scan stat tiles, BRS trend, system pillars")

    set_notes(s, "KAVACH's answer, at a glance: repository in, CI/CD trigger, nine scanners in parallel, "
                 "AI explanation layer, a business-aware risk engine, deterministic compliance mapping, and "
                 "seven report formats out. Every box on this flow corresponds to a real service in the "
                 "codebase -- we'll open each one up over the next several slides. The screenshot on the "
                 "right is the authenticated overview dashboard a security engineer sees first.")


# ------------------------------------------------------------------ 06 ----
def slide_06_comparison(prs):
    s = add_slide(prs, "Competitive Landscape", 6, TOTAL)
    kicker_title(s, "Competitive Landscape", "Why Existing Tools Fall Short For Banking", title_size=28)

    headers = ["Capability", "SonarQube", "Snyk", "GitHub Adv. Security", "OWASP Dep-Check", "Semgrep (OSS)", "KAVACH"]
    rows = [
        ["Static + dependency + config + secrets, one pipeline", "Partial", "Partial", "Partial", "Dep. only", "Code only", "Yes -- 9 tools, one scan"],
        ["Business-aware risk score (not flat CVSS)", "No", "No", "No", "No", "No", "Banking Risk Score, 7 factors"],
        ["RBI / PCI DSS / SWIFT CSP clause mapping", "No", "No", "No", "No", "No", "Yes -- YAML-driven, live"],
        ["Cross-tool dedup & correlation", "N/A", "N/A", "N/A", "N/A", "N/A", "Yes -- CVE/file/category match"],
        ["Grounded, citation-backed AI explanation", "No", "Basic", "Basic (Copilot Autofix)", "No", "No", "RAG + confidence-gated"],
    ]
    themed_table(s, Inches(LEFT), Inches(2.05), Inches(CW), Inches(3.75), headers, rows,
                 col_widths=[Inches(3.05), Inches(1.33), Inches(1.12), Inches(1.59), Inches(1.44), Inches(1.38), Inches(2.32)],
                 highlight_col=6, body_size=10.5)

    source_note(s, Inches(LEFT), Inches(5.95), Inches(CW),
                "Comparison reflects each tool's typical, publicly-documented scope as of this writing -- "
                "not a claim that these tools lack all capability, but that banking-specific risk scoring "
                "and regulatory mapping are not their focus. KAVACH's compliance mapping is explicitly its "
                "own illustrative rule set (see Slide 14), not a certified audit.")

    set_notes(s, "We're not claiming SonarQube, Snyk, GitHub Advanced Security, OWASP Dependency-Check, or "
                 "Semgrep are bad tools -- several of them, including Semgrep itself, are running inside "
                 "KAVACH. The honest differentiation is: none of them ship a banking-specific business risk "
                 "score, none of them auto-map findings to RBI/PCI DSS/SWIFT CSP clauses, and none of them "
                 "run this specific 9-tool cross-correlation with a citation-gated AI layer on top. That's "
                 "KAVACH's lane: not a better Semgrep, but the business and regulatory layer none of them "
                 "provide.")


# ------------------------------------------------------------------ 07 ----
def slide_07_architecture(prs):
    s = add_slide(prs, "System Architecture", 7, TOTAL)
    kicker_title(s, "System Architecture", "22 Components, 6 Layers -- All Implemented", title_size=27)

    def row(y, items, node_h=0.5, max_w=2.7, size=9.5):
        n = len(items)
        w = min(max_w, (CW - 0.12 * (n - 1)) / n)
        total = w * n + 0.12 * (n - 1)
        x0 = LEFT + (CW - total) / 2
        for i, (label, color) in enumerate(items):
            x = x0 + i * (w + 0.12)
            arch_node(s, Inches(x), Inches(y), Inches(w), Inches(node_h), label, accent=color, size=size)
        return x0 + total / 2

    def tag(y, text, color=GOLD):
        textbox(s, Inches(LEFT), Inches(y), Inches(4), Inches(0.18), text, size=8.5, color=color,
                bold=True, letter_spacing=0.8)

    def arrow(y):
        connector(s, Inches(LEFT + CW / 2), Inches(y), Inches(LEFT + CW / 2), Inches(y + 0.11), color=BORDER_LT)

    y = 2.02
    tag(y, "INGRESS", BLUE_LT)
    row(y + 0.19, [("Repository", BLUE), ("Webhook", BLUE), ("Gateway", BLUE), ("Authentication", BLUE)])
    arrow(y + 0.19 + 0.5)
    y += 0.19 + 0.5 + 0.12
    tag(y, "ORCHESTRATION", RED_LT)
    row(y + 0.19, [("Queue (Redis / Celery, priority)", RED), ("Distributed Workers (fan-out)", RED)], max_w=4.0)
    arrow(y + 0.19 + 0.5)
    y += 0.19 + 0.5 + 0.12
    tag(y, "PARALLEL SCANNERS -- 9 INDEPENDENT TOOLS", BLUE_LT)
    row(y + 0.19, [("Semgrep", BLUE_LT), ("Joern", BLUE_LT), ("AST-Grep", BLUE_LT),
                   ("Dependency Scanner", BLUE_LT), ("Configuration Scanner", BLUE_LT), ("Secrets Detection", BLUE_LT)],
        node_h=0.46, size=8.7)
    arrow(y + 0.19 + 0.46)
    y += 0.19 + 0.46 + 0.12
    tag(y, "INTELLIGENCE", GOLD)
    row(y + 0.19, [("Aggregation Layer", SLATE), ("Banking Risk Score Engine", GOLD),
                   ("Compliance Engine", GREEN), ("AI Explanation Layer", PURPLE_LT)])
    arrow(y + 0.19 + 0.5)
    y += 0.19 + 0.5 + 0.12
    tag(y, "KNOWLEDGE (RAG)", PURPLE_LT)
    row(y + 0.19, [("Knowledge Base (pgvector)", PURPLE), ("Executive Intelligence", PURPLE)], max_w=4.2)
    arrow(y + 0.19 + 0.5)
    y += 0.19 + 0.5 + 0.12
    tag(y, "DELIVERY", GREEN_LT)
    row(y + 0.19, [("Report Generator", BLUE), ("Dashboard", BLUE), ("Notifications", GOLD), ("Storage (Local/S3)", SLATE)])

    set_notes(s, "This is the full component inventory from the README's own architecture diagram, laid out "
                 "as six layers: ingress (repository, webhook, gateway, auth), orchestration (Redis-backed "
                 "priority queue and distributed Celery workers), nine parallel scanners, an intelligence "
                 "layer (aggregation, Banking Risk Score, compliance, AI explanation), the RAG knowledge "
                 "layer, and delivery (reports, dashboard, notifications, storage). Every node on this "
                 "slide is a real module in backend/app/ -- we cited file paths for each during our own "
                 "internal audit.")


# ------------------------------------------------------------------ 08 ----
def slide_08_pipeline(prs):
    s = add_slide(prs, "Scan Pipeline", 8, TOTAL)
    kicker_title(s, "Scan Pipeline", "From Push to Dashboard, Automatically")

    steps = ["Repository", "Webhook", "Queue", "Parallel\nScanners", "Aggregation", "AI", "Compliance", "Reports", "Dashboard"]
    process_flow(s, steps, Inches(LEFT), Inches(2.5), Inches(CW), Inches(1.05),
                 colors=[BLUE, BLUE, RED, BLUE_LT, SLATE, PURPLE, GREEN, GOLD, WHITE], label_size=10.5)

    icon_card(s, Inches(LEFT), Inches(4.0), Inches(3.95), Inches(2.65), "Entry Points",
              "A scan starts from a user-submitted .zip upload, a direct repository URL, a bundled sandbox "
              "payload, a nightly scheduled rescan, or a verified GitHub push webhook (HMAC-checked via "
              "X-Hub-Signature-256) -- all converge on the same orchestration path.", accent=BLUE)
    icon_card(s, Inches(4.14), Inches(4.0), Inches(3.95), Inches(2.65), "Fault-Isolated Fan-Out",
              "prepare_scan_job dispatches a Celery chord: all 9 scanner tasks run in parallel, each with "
              "its own timeout and 3-attempt retry with backoff. A scanner failing (or a missing local "
              "tool) never blocks the others -- the chord callback always fires.", accent=RED)
    icon_card(s, Inches(8.28), Inches(4.0), Inches(3.95), Inches(2.65), "Live Progress",
              "Per-scanner status lives in Redis (queued/running/completed/failed) and streams to the "
              "client over a WebSocket (/scan/{id}/ws) -- the dashboard shows real-time progress, not a "
              "spinner.", accent=GOLD)

    set_notes(s, "A scan can enter from three places -- manual upload, a direct URL, or a verified GitHub "
                 "webhook -- and all three converge on one orchestration path. The queue dispatches a "
                 "Celery chord of all nine scanners in parallel; the chord callback -- aggregation -- fires "
                 "automatically once every scanner returns, success or failure, because each scanner task "
                 "catches its own exceptions and always returns a result. Progress streams live over "
                 "WebSocket, backed by Redis, so the dashboard reflects real per-scanner state, not a "
                 "fake progress bar.")


# ------------------------------------------------------------------ 09 ----
def slide_09_scanners(prs):
    s = add_slide(prs, "Parallel Scanning Engine", 9, TOTAL)
    kicker_title(s, "Parallel Scanning Engine", "9 Independent Celery Tasks, Fanned Out Per Scan", title_size=26)

    scanners = [
        ("Semgrep", "Custom rule pack: secrets, SQLi/cmd injection, weak crypto, unsafe deserialization, "
                     "path traversal, insecure randomness. Regex fallback if the binary is unavailable.", BLUE_LT),
        ("Joern", "Code-property-graph reachability from dangerous sinks (exec, eval, pickle.loads) back to "
                   "untrusted input -- runs when a local Joern install is present.", BLUE_LT),
        ("AST-Grep", "Structural AST rules (eval/exec, unsafe child_process.exec, asserts used as security "
                      "checks) -- deliberately independent of Semgrep's pattern language.", BLUE_LT),
        ("Dependency Scanner", "pip-audit + a CycloneDX SBOM, cross-checked against OSV.dev batch queries "
                                 "and a rate-limited NVD keyword search -- 3 independent CVE lookups.", GOLD),
        ("Configuration Scanner", "Structural YAML analysis of Kubernetes manifests, Docker Compose, and "
                                    "GitHub Actions, plus a dedicated Dockerfile instruction analyzer.", GREEN_LT),
        ("Secrets Detection", "In-house, gitleaks-style regex engine: AWS/GitHub/GitLab/Slack/Stripe keys, "
                                 "private-key blocks, JWTs, credentials embedded in URLs.", RED_LT),
    ]
    w = (CW - 0.16 * 2) / 3
    for i, (name, desc, color) in enumerate(scanners):
        x = LEFT + (i % 3) * (w + 0.16)
        y = 2.05 + (i // 3) * 1.62
        icon_card(s, Inches(x), Inches(y), Inches(w), Inches(1.5), name, desc, accent=color, title_size=12.5, body_size=9.7)

    card(s, Inches(LEFT), Inches(5.35), Inches(CW), Inches(1.35), fill=PANEL_ALT, line=BORDER)
    textbox(s, Inches(LEFT + 0.25), Inches(5.48), Inches(CW - 0.5), Inches(0.3),
            "PRIORITY-QUEUED, FAULT-ISOLATED EXECUTION", size=11, color=GOLD, bold=True, letter_spacing=1.0)
    textbox(s, Inches(LEFT + 0.25), Inches(5.8), Inches(CW - 0.5), Inches(0.8),
            "Named Redis queues -- kavach.critical, kavach.high, kavach.normal, kavach.low -- back a "
            "dedicated always-on worker pool for critical/high scans plus a horizontally-scalable pool for "
            "bulk throughput. Every scanner task retries up to 3 times with exponential backoff and always "
            "returns a result, so one missing tool degrades gracefully instead of blocking the scan.",
            size=11.5, color=SLATE, line_spacing=1.2)

    set_notes(s, "Nine independent Celery tasks, not nine steps of one script: Semgrep, Joern, AST-Grep, a "
                 "dependency scanner that itself runs pip-audit plus OSV.dev plus an NVD keyword search, a "
                 "configuration scanner for Kubernetes/Compose/GitHub Actions plus Dockerfiles, and a "
                 "dedicated secrets engine. They run on named priority queues -- kavach.critical through "
                 "kavach.low -- with a dedicated worker pool for urgent scans so a backlog of low-priority "
                 "jobs never starves a critical one. Every task retries independently and always returns, "
                 "so a missing local tool -- say, no Joern install -- degrades gracefully rather than "
                 "failing the whole scan.")


# ------------------------------------------------------------------ 10 ----
def slide_10_ase(prs):
    s = add_slide(prs, "Attack Surface Exposure", 10, TOTAL)
    kicker_title(s, "Attack Surface Exposure", "A Heuristic Exposure Index -- Not a Zero-Day Prediction", title_size=25)

    card(s, Inches(LEFT), Inches(2.05), Inches(7.1), Inches(1.0), fill=PANEL_ALT, line=GOLD, line_w=1.1)
    textbox(s, Inches(LEFT + 0.25), Inches(2.2), Inches(6.6), Inches(0.75),
            "We deliberately do not call this “Zero-Day Prediction.” It is a composite index of "
            "factors correlated with exposure to undisclosed vulnerabilities -- not a forecast of any "
            "specific future exploit.", size=13, color=WHITE, italic=True, line_spacing=1.15)

    factors = [
        ("Dependency count", "0.18 × SBOM component count, capped 20"),
        ("Known CVE density", "4.0 × critical/high pip-audit findings, capped 30"),
        ("Dependency staleness", "4.0 × est. avg. dependency age (yrs), capped 20"),
        ("Risky package categories", "1.5 × crypto/auth/serialization-adjacent packages, capped 15"),
        ("Configuration risk", "critical×3.0 + high×1.5 config findings, capped 10"),
        ("Code vulnerability density", "2.5 × critical static findings, capped 15"),
    ]
    ty = 3.2
    for label, formula in factors:
        card(s, Inches(LEFT), Inches(ty), Inches(7.1), Inches(0.55), fill=PANEL, line=BORDER)
        textbox(s, Inches(LEFT + 0.2), Inches(ty + 0.06), Inches(2.6), Inches(0.42), label, size=11.5,
                color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, Inches(LEFT + 2.85), Inches(ty + 0.06), Inches(4.1), Inches(0.42), formula, size=10,
                color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
        ty += 0.61

    icon_card(s, Inches(8.0), Inches(2.05), Inches(4.23), Inches(1.55), "How It's Used",
              "Sum of 6 capped factors, capped again at 100. Bands: Critical ≥70, High ≥45, "
              "Medium ≥20, Low <20. Persisted independently of BRS on its own ScanResult columns.",
              accent=BLUE)
    icon_card(s, Inches(8.0), Inches(3.7), Inches(4.23), Inches(1.55), "Its Purpose",
              "Flags dependency-rot and risky-surface accumulation even on a scan with zero confirmed "
              "findings -- a codebase can be BRS-clean today and still be sitting on this problem.",
              accent=GOLD)
    icon_card(s, Inches(8.0), Inches(5.35), Inches(4.23), Inches(1.35), "Explicit Limitations",
              "Fixed confidence = 0.55, labeled a prototype in-code. Documented path to replace the "
              "heuristic with a calibrated classifier -- not shipped as one today.", accent=RED)

    set_notes(s, "We want to be precise about terminology here, because an earlier version of our own "
                 "materials called this 'Zero-Day Prediction' -- the current code renamed it explicitly, "
                 "because that name overclaimed what a heuristic can support. Attack Surface Exposure is a "
                 "0-100 composite of six capped, additive factors -- dependency count, known-CVE density, "
                 "dependency staleness, risky package categories, configuration risk, and code vulnerability "
                 "density. It's independent of BRS by design -- a codebase can score a clean BRS today and "
                 "still carry a high Attack Surface Exposure from dependency rot. It ships with a fixed "
                 "confidence of 0.55 and is explicitly documented in code as a prototype, with a stated path "
                 "to a calibrated classifier later.")


# ------------------------------------------------------------------ 11 ----
def slide_11_brs(prs):
    s = add_slide(prs, "Banking Risk Score", 11, TOTAL)
    kicker_title(s, "Business-Aware Risk Scoring", "Technical Severity + Business Context = Banking Risk Score", title_size=25)

    process_flow(s, ["Technical\nSeverity", "+", "Business\nContext", "=", "Banking\nRisk Score"],
                 Inches(LEFT), Inches(2.02), Inches(7.1), Inches(0.7),
                 colors=[BLUE, PANEL, GOLD, PANEL, GREEN], label_size=11)

    factors = [
        ("CVSS", "30%"), ("Business criticality", "20%"), ("Exploitability", "15%"),
        ("Internet exposure", "10%"), ("Compliance impact", "10%"), ("Asset value", "10%"),
        ("Historical incidents", "5%"),
    ]
    ty = 2.95
    for label, pct in factors:
        card(s, Inches(LEFT), Inches(ty), Inches(7.1), Inches(0.44), fill=PANEL, line=BORDER)
        bar_w = 7.1 * 0.62 * (float(pct.strip('%')) / 30.0)
        rect(s, Inches(LEFT), Inches(ty), Inches(min(bar_w, 7.1)), Inches(0.44), fill=GOLD, line=None, radius=0.35)
        textbox(s, Inches(LEFT + 0.18), Inches(ty + 0.03), Inches(4.2), Inches(0.38), label, size=11.5,
                color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, Inches(LEFT + 6.35), Inches(ty + 0.03), Inches(0.65), Inches(0.38), pct, size=12,
                color=GOLD_LT, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        ty += 0.485

    icon_card(s, Inches(8.0), Inches(2.02), Inches(4.23), Inches(1.95), "Risk Bands (0-100)",
              "Critical ≥82  ·  High 58-82  ·  Medium 35-58  ·  Low <35 -- empirically "
              "calibrated and unit-tested at every boundary.", accent=RED)
    icon_card(s, Inches(8.0), Inches(4.05), Inches(4.23), Inches(1.75), "Business Modules",
              "Payments (10/10, internet-facing), Authentication (8.5/8.0), Customer Data (7.0/9.0), "
              "Admin, Infrastructure, Reporting, General -- criticality × asset value, both DB-editable.",
              accent=GOLD)
    icon_card(s, Inches(8.0), Inches(5.88), Inches(4.23), Inches(0.9), "Runtime-Configurable",
              "Every weight and module lives in Postgres, editable via /api/v1/risk/* -- no redeploy.",
              accent=GREEN, body_size=9.5)

    set_notes(s, "The Banking Risk Score is a weighted average across seven factors -- CVSS at 30%, "
                 "business-module criticality at 20%, exploitability at 15%, internet exposure and "
                 "compliance impact and asset value at 10% each, and historical incident count at 5% -- "
                 "scaled to 0-100 and clamped. It's a weighted average deliberately, not a product, because "
                 "multiplying seven normalized terms collapses toward zero unless every factor is high, "
                 "which would hide a genuinely severe single factor. Risk bands are Critical at 82+, High "
                 "58-82, Medium 35-58, Low below 35 -- calibrated against real scan data and pinned down "
                 "with boundary unit tests. Every weight and every business module -- Payments, "
                 "Authentication, Customer Data, and so on -- lives in Postgres and is editable at runtime "
                 "through the risk-config API, with no redeploy required.")


# ------------------------------------------------------------------ 12 ----
def slide_12_ai(prs):
    s = add_slide(prs, "AI Intelligence Layer", 12, TOTAL)
    kicker_title(s, "AI Intelligence Layer", "Strictly Explanatory -- Never the Scorer")

    icon_card(s, Inches(LEFT), Inches(2.05), Inches(5.85), Inches(2.15), "1. Automated Finding Explanations",
              "Pipeline-embedded, no user interaction needed. A provider-agnostic gateway resolves Claude, "
              "OpenAI, and Gemini (cloud) or Ollama and vLLM (local) behind one AI_MODE setting -- hybrid "
              "by default: local first, cloud fallback. A semantic cache avoids re-explaining near-identical "
              "findings, and a rule-based template library guarantees an explanation even with zero "
              "providers configured -- a scan never blocks on AI availability.", accent=PURPLE, body_size=10.3)
    icon_card(s, Inches(6.13), Inches(2.05), Inches(6.1), Inches(2.15), "2. RAG Knowledge Layer",
              "On-demand, citation-grounded: AI Assistant chat, per-Finding Intelligence, and Executive "
              "Intelligence share one retrieve → rerank → confidence-gate → generate "
              "pipeline against an uploaded knowledge base (see next slide).", accent=PURPLE_LT, body_size=10.5)

    card(s, Inches(LEFT), Inches(4.4), Inches(CW), Inches(1.35), fill=PANEL_ALT, line=GOLD, line_w=1.25)
    textbox(s, Inches(LEFT + 0.3), Inches(4.55), Inches(CW - 0.6), Inches(0.32),
            "A STRUCTURAL GUARANTEE, NOT A POLICY", size=11.5, color=GOLD, bold=True, letter_spacing=1.0)
    textbox(s, Inches(LEFT + 0.3), Inches(4.88), Inches(CW - 0.6), Inches(0.8),
            "“AI never calculates a security score. AI never changes a finding. AI only explains, "
            "summarizes, retrieves evidence, and assists a human who is still the one deciding what to "
            "do.” The AI service layer only ever reads Finding/ScanJob rows to build a prompt -- none "
            "of it commits back to the database, and BRS/Attack-Surface-Exposure scoring runs in a "
            "separate code path the AI layer never calls into.", size=12.5, color=WHITE, italic=True, line_spacing=1.18)

    screenshot_placeholder(s, Inches(LEFT), Inches(5.9), Inches(3.9), Inches(0.95), "Executive Summary",
                            route="/executive", note="Executive Intelligence panel")
    mini_stat(s, Inches(4.6), Inches(5.9), Inches(3.75), Inches(0.95), "5 providers", "Claude · OpenAI · "
              "Gemini · Ollama · vLLM, one gateway", accent=BLUE)
    mini_stat(s, Inches(8.5), Inches(5.9), Inches(3.73), Inches(0.95), "0 hard fails", "Template fallback -- "
              "scans never block on AI reachability", accent=GREEN)

    set_notes(s, "Two AI subsystems, addressing two different needs. First: every finding gets an automatic "
                 "explanation as part of the scan pipeline, from a provider-agnostic gateway that can call "
                 "Claude, OpenAI, or Gemini in the cloud, or Ollama and vLLM locally, behind one hybrid-mode "
                 "setting that tries local first and falls back to cloud. If no provider is configured at "
                 "all, a rule-based template library still produces a deterministic explanation -- a scan's "
                 "completion is never gated on AI availability. Second: an on-demand retrieval-augmented "
                 "layer covering chat, per-finding deep-dives, and executive Q&A, detailed next. Structurally "
                 "-- not just as a policy -- the AI layer only ever reads finding data to build a prompt; it "
                 "never writes a score or mutates a finding. BRS and Attack Surface Exposure run in a "
                 "completely separate code path the AI layer never touches.")


# ------------------------------------------------------------------ 13 ----
def slide_13_rag(prs):
    s = add_slide(prs, "Knowledge Base (RAG)", 13, TOTAL)
    kicker_title(s, "Knowledge Base", "Retrieval-Augmented, Confidence-Gated, Always Cited", title_size=26)

    steps = ["Ingest", "Embed", "Vector DB", "Retrieve\n(top 20)", "Rerank\n(top 5)", "Confidence\nGate", "Generate /\nRefuse", "Cite"]
    process_flow(s, steps, Inches(LEFT), Inches(2.05), Inches(CW), Inches(0.95),
                 colors=[BLUE, PURPLE, PURPLE, BLUE_LT, BLUE_LT, GOLD, GREEN, WHITE], label_size=9.7)

    details = [
        ("Ingestion", "PDF / Markdown / text, deduplicated by SHA-256 content hash, version-chained by "
                        "filename, chunked with heading/section/page-aware splitting (~350 tokens, 50 overlap)."),
        ("Embedding + Retrieval", "BAAI/bge-small-en-v1.5 (384-dim, local ONNX, no GPU) into Postgres + "
                                    "pgvector HNSW cosine search -- top 20 candidates retrieved."),
        ("Rerank + Confidence Gate", "Xenova/ms-marco-MiniLM-L-6-v2 cross-encoder narrows to top 5; the "
                                        "sigmoid-normalized top score must clear 0.5 or the LLM is never "
                                        "called at all."),
        ("Citations", "Every answer carries document, section, page, cosine similarity and rerank score -- "
                        "not a vague “source: internet.”"),
    ]
    w = (CW - 0.16 * 3) / 4
    for i, (t, d) in enumerate(details):
        x = LEFT + i * (w + 0.16)
        icon_card(s, Inches(x), Inches(3.2), Inches(w), Inches(2.0), t, d, accent=[BLUE, PURPLE, GOLD, GREEN][i],
                   title_size=11.5, body_size=9.3)

    card(s, Inches(LEFT), Inches(5.4), Inches(CW), Inches(1.3), fill=PANEL_ALT, line=GOLD, line_w=1.1)
    textbox(s, Inches(LEFT + 0.3), Inches(5.53), Inches(CW - 0.6), Inches(0.3),
            "WHY THIS IS SAFER THAN PLAIN LLM PROMPTING", size=11, color=GOLD, bold=True, letter_spacing=1.0)
    textbox(s, Inches(LEFT + 0.3), Inches(5.85), Inches(CW - 0.6), Inches(0.78),
            "Below the confidence threshold, KAVACH physically never calls the model -- it returns a fixed "
            "refusal message deterministically. Above threshold, the LLM sees only the retrieved excerpts "
            "(closed-context generation), so it cannot answer from ungrounded general knowledge. Every "
            "surface is citation-backed by construction, not by prompt instruction alone.",
            size=11, color=SLATE, line_spacing=1.15)

    set_notes(s, "The knowledge base pipeline: an administrator uploads PDF, Markdown, or text reference "
                 "material -- OWASP guides, internal policy, regulatory text -- deduplicated by content "
                 "hash and chunked with heading and page awareness. Chunks are embedded locally with "
                 "BAAI/bge-small-en-v1.5, no external API call, into Postgres with pgvector's HNSW index. A "
                 "query retrieves the top 20 chunks by cosine similarity, a cross-encoder reranks down to "
                 "the top 5, and a sigmoid-normalized confidence score must clear 0.5 before the LLM is even "
                 "invoked. Below that threshold, KAVACH returns a fixed refusal message -- deterministically, "
                 "not as a prompt suggestion the model could ignore. Every answer that does generate carries "
                 "real citations: document, section, page, and both similarity scores. That's the core "
                 "safety argument over plain prompting -- the model can't hallucinate an answer from its own "
                 "training data, because it never sees the question without the gate first deciding whether "
                 "grounding evidence actually exists.")


# ------------------------------------------------------------------ 14 ----
def slide_14_compliance(prs):
    s = add_slide(prs, "Compliance Engine", 14, TOTAL)
    kicker_title(s, "Compliance Engine", "Deterministic, YAML-Driven Mapping -- Live, Not Frozen", title_size=25)

    fws = [
        ("RBI IT Framework, 2021", GREEN, ["4.2 Access Control & Identity Mgmt", "5.3 Cryptography Policy",
                                             "6.4 Secure Coding Standards", "6.6 Patch & Vulnerability Mgmt"]),
        ("PCI DSS v4.0", BLUE, ["2.2 Secure configurations", "6.2 Bespoke software security",
                                  "8.3 Strong authentication", "12.8 Third-party risk"]),
        ("SWIFT CSP", PURPLE_LT, ["2.6 Cryptographic Controls", "2.7 Vulnerability Scanning",
                                     "3.1 Software Integrity", "6.1 Operator Session Security"]),
    ]
    w = (CW - 0.2 * 2) / 3
    for i, (name, color, ctrls) in enumerate(fws):
        x = LEFT + i * (w + 0.2)
        card(s, Inches(x), Inches(2.05), Inches(w), Inches(2.55), fill=PANEL, line=color, line_w=1.25)
        textbox(s, Inches(x + 0.22), Inches(2.2), Inches(w - 0.44), Inches(0.4), name, size=13.5, color=WHITE, bold=True)
        yy = 2.68
        for c in ctrls:
            dot(s, Inches(x + 0.28), Inches(yy + 0.08), Inches(0.1), color)
            textbox(s, Inches(x + 0.46), Inches(yy), Inches(w - 0.65), Inches(0.4), c, size=9.7, color=SLATE, line_spacing=1.05)
            yy += 0.44

    icon_card(s, Inches(LEFT), Inches(4.78), Inches(3.95), Inches(1.9), "How It Evaluates",
              "Each YAML control defines a trigger (category + min-severity + source + keywords). Any "
              "matching finding fails the control, with full evidence attached -- binary PASS/FAIL, no "
              "middle state.", accent=GOLD, body_size=10)
    icon_card(s, Inches(4.14), Inches(4.78), Inches(3.95), Inches(1.9), "Data-Driven, Not Hardcoded",
              "Dropping a new .yaml file into app/data/compliance_rules/ is the entire integration step for "
              "a new framework -- zero code changes. Recomputed live per scan, on demand.", accent=BLUE, body_size=10)
    icon_card(s, Inches(8.28), Inches(4.78), Inches(3.95), Inches(1.9), "Honest Scope",
              "KAVACH's own illustrative mapping for continuous self-assessment -- not a certified PCI QSA, "
              "RBI, or SWIFT CSP attestation. Should not be presented to a regulator as one without "
              "independent review.", accent=RED, body_size=10)

    set_notes(s, "Compliance evaluation is entirely YAML-driven -- every control ID, title, and trigger "
                 "condition for RBI IT Framework 2021, PCI DSS v4.0, and SWIFT CSP lives in a YAML file, not "
                 "in Python. A trigger matches on finding category, minimum severity, source tool, and "
                 "keywords; any match fails that control, with evidence -- file, line, severity, source -- "
                 "attached. It's a binary PASS/FAIL per control, rolled up into a per-framework and overall "
                 "compliance percentage, recomputed live on every request rather than frozen at scan time. "
                 "We want to be very direct about scope here: this is KAVACH's own illustrative mapping, "
                 "useful for continuous self-assessment and evidence-gathering -- it is explicitly not a "
                 "certified PCI QSA, RBI, or SWIFT CSP attestation, and we say that in the code comments "
                 "themselves, not just in this deck.")


# ------------------------------------------------------------------ 15 ----
def slide_15_rbac(prs):
    s = add_slide(prs, "Role-Based Access Control", 15, TOTAL)
    kicker_title(s, "Role-Based Access Control", "5 Backend-Enforced Roles, 13 Permissions", title_size=27)

    headers = ["Role", "Display Name", "Perms", "Typical Use"]
    rows = [
        ["admin", "Administrator", "13 / 13", "User & role mgmt, risk config, full scan & report access"],
        ["security_engineer", "Security Manager", "11", "Runs/triages scans, risk config, audit log, team analytics"],
        ["developer", "Security Analyst", "9", "Submits scans, findings, compliance, AI Assistant, KB write"],
        ["auditor", "Executive / Board Member", "7", "Read-only risk, compliance, executive reporting, audit log"],
        ["read_only", "Read Only", "3", "Self-registration default -- scan/report visibility only"],
    ]
    themed_table(s, Inches(LEFT), Inches(2.05), Inches(CW), Inches(2.85), headers, rows,
                 col_widths=[Inches(2.3), Inches(2.9), Inches(1.3), Inches(5.73)], highlight_col=None, body_size=11.5)

    icon_card(s, Inches(LEFT), Inches(5.15), Inches(3.95), Inches(1.55), "Enforcement",
              "Coarse middleware blocks any mutating request (POST/PUT/PATCH/DELETE) from read-only-shaped "
              "roles outright; fine-grained per-route Permission checks enforce the rest.", accent=GOLD, body_size=9.7)
    icon_card(s, Inches(4.14), Inches(5.15), Inches(3.95), Inches(1.55), "Not Just UX Gating",
              "Web and mobile clients read the same role table only to hide navigation -- it's a UX "
              "convenience. The backend is the real, only security boundary.", accent=BLUE, body_size=9.7)
    icon_card(s, Inches(8.28), Inches(5.15), Inches(3.95), Inches(1.55), "Auditable",
              "Every login attempt, role change, and permission denial is persisted with actor, IP, and "
              "outcome -- queryable via a dedicated audit-log endpoint.", accent=GREEN, body_size=9.7)

    set_notes(s, "Five roles, enforced at the backend two ways: a coarse middleware that blocks any "
                 "mutating HTTP verb from a read-only-shaped role outright, and fine-grained per-route "
                 "permission checks across 13 distinct permission strings. Administrator holds all 13; "
                 "Security Manager runs and triages scans with risk-config write access; Security Analyst "
                 "submits scans and uses the AI Assistant; Executive / Board Member is deliberately "
                 "read-only across risk, compliance, and reporting; Read Only is the self-registration "
                 "default. Both the web and mobile clients read this same role table, but only to hide "
                 "navigation -- that's UX convenience, never the actual security boundary. Every login, role "
                 "change, and permission denial is audit-logged with actor, IP address, and outcome.")


# ------------------------------------------------------------------ 16 ----
def slide_16_mobile(prs):
    s = add_slide(prs, "Flutter Mobile Application", 16, TOTAL)
    kicker_title(s, "Flutter Mobile Application", "Same Backend, Same RBAC, Native Android/iOS", title_size=26)

    icon_card(s, Inches(LEFT), Inches(2.05), Inches(4.0), Inches(2.15), "Architecture",
              "Riverpod state management, Dio + JWT refresh interceptor mirroring the web client, "
              "go_router with backend-role-aware redirects, and Freezed/json_serializable models mirroring "
              "backend Pydantic schemas field-for-field. Owns no business logic of its own.", accent=BLUE, body_size=10.3)

    textbox(s, Inches(LEFT), Inches(4.35), Inches(4.0), Inches(0.3), "FULLY WIRED, REAL ENDPOINTS", size=11,
            color=GREEN, bold=True, letter_spacing=0.8)
    bullets(s, Inches(LEFT), Inches(4.68), Inches(4.0), Inches(2.1), [
        "Login / Signup / Session restore",
        "Dashboard -- my-activity analytics, severity breakdown",
        "Repositories -- list, scheduled-scan toggle",
        "Start Scan -- zip upload or URL, priority selector",
        "Scan Queue & Details -- live status, cancel action",
    ], size=10.3, marker_color=GREEN, gap=0.09)

    textbox(s, Inches(4.28), Inches(2.05), Inches(3.7), Inches(0.3), "HONESTLY-LABELED PLACEHOLDERS", size=11,
            color=GOLD, bold=True, letter_spacing=0.8)
    bullets(s, Inches(4.28), Inches(2.38), Inches(3.7), Inches(2.5), [
        ("Risk Dashboard, Finding Explorer, Compliance -- ", "no cross-scan/portfolio rollup endpoint "
                                                                 "exists yet on the backend."),
        ("Executive Summary -- ", "RAG-based endpoint out of scope for this milestone."),
        ("Notifications, Settings -- ", "no notifications API or self-service profile endpoint exists yet."),
    ], size=10, marker_color=GOLD, gap=0.12)
    source_note(s, Inches(4.28), Inches(4.95), Inches(3.7),
                "Disclosed explicitly in mobile/docs/backend_gaps.md -- placeholders by design, not fake data.")

    screenshot_placeholder(s, Inches(8.3), Inches(2.05), Inches(3.93), Inches(4.65), "Flutter Mobile App",
                            route="Dashboard / Repositories / Scans", note="Android/iOS, shared backend")

    set_notes(s, "The Flutter app is a real, native Android/iOS client -- Riverpod for state, a Dio client "
                 "with a JWT refresh interceptor mirroring the web app's, go_router with the same "
                 "backend-role-aware redirects, and Freezed models generated to match backend Pydantic "
                 "schemas field-for-field. It owns no business logic of its own -- every screen calls the "
                 "identical FastAPI backend. Dashboard, Repositories, Start Scan, and Scan Queue/Details are "
                 "fully wired to real endpoints today. In the interest of not fabricating anything: Risk "
                 "Dashboard, Finding Explorer, Compliance, Executive Summary, Notifications, and Settings "
                 "are honestly-labeled placeholders, because the backend doesn't yet expose the "
                 "cross-scan rollup or notifications endpoints they'd need -- this is documented directly in "
                 "the repo's mobile/docs/backend_gaps.md, not glossed over.")


# ------------------------------------------------------------------ 17 ----
def slide_17_frontend(prs):
    s = add_slide(prs, "Frontend Experience", 17, TOTAL)
    kicker_title(s, "Frontend Experience", "React 19 + Tailwind v4 -- A Genuine Two-Palette Dark Mode", title_size=25)

    bullets(s, Inches(LEFT), Inches(2.05), Inches(4.55), Inches(3.1), [
        ("13+ role-aware pages -- ", "Risk, Compliance, Finding Explorer, Executive, Knowledge Base, AI "
                                        "Assistant, RAG Ops, Admin, and more."),
        ("Real-time scan lifecycle -- ", "WebSocket-driven live per-scanner progress, not polling."),
        ("Framer Motion animation ", "throughout; Recharts for BRS trend, severity, and compliance charts."),
        ("True dual-palette dark mode -- ", "a real ThemeContext driving both Tailwind classes and "
                                                "inline chart colors, not default scaffolding."),
    ], size=12, gap=0.16)

    icon_card(s, Inches(LEFT), Inches(5.15), Inches(4.55), Inches(1.65), "The 3D Architecture Explorer",
              "React Three Fiber scene: adaptive quality tiers (mobile/tablet/desktop), hover tooltips, "
              "click-to-focus glassmorphism side panels, GSAP camera flights, and 13 custom per-node "
              "geometries -- available publicly (no login) and inside the authenticated dashboard.",
              accent=PURPLE_LT, body_size=9.8)

    screenshot_placeholder(s, Inches(4.85), Inches(2.05), Inches(3.68), Inches(4.75), "Architecture Explorer",
                            route="/architecture (public) or /dashboard/architecture",
                            note="3D, interactive, click any node")
    screenshot_placeholder(s, Inches(8.68), Inches(2.05), Inches(3.55), Inches(4.75), "Scan Results",
                            route="/scans/:scanJobId", note="Live status + finding detail modal")

    set_notes(s, "The web console is the deep-work surface: thirteen-plus role-aware pages, a WebSocket-driven "
                 "live scan lifecycle instead of polling, Framer Motion animation throughout, and Recharts "
                 "visualizations that actually respond to the theme -- dark mode here is a genuine two-palette "
                 "system with its own color tokens, not default Tailwind scaffolding. The showcase feature is "
                 "the 3D architecture explorer built on React Three Fiber: it adapts rendering cost to "
                 "device tier, gives every node a distinct custom geometry, and opens a glassmorphism detail "
                 "panel on click with GSAP-animated camera flights. It's available publicly, with no login, "
                 "specifically so evaluators and recruiters can explore the system design without needing an "
                 "account.")


# ------------------------------------------------------------------ 18 ----
def slide_18_cicd(prs):
    s = add_slide(prs, "CI/CD Integration", 18, TOTAL)
    kicker_title(s, "CI/CD Integration", "Webhook-Triggered Scanning, Straight Into Existing Tooling", title_size=25)

    process_flow(s, ["Git Push", "Webhook", "Pipeline", "Analysis", "Reports", "Developer\nFeedback"],
                 Inches(LEFT), Inches(2.15), Inches(CW), Inches(0.95),
                 colors=[BLUE, BLUE, RED, BLUE_LT, GOLD, GREEN], label_size=11.5)

    icon_card(s, Inches(LEFT), Inches(3.45), Inches(3.95), Inches(2.2), "KAVACH's Own Pipeline",
              "GitHub Actions: backend-test (pytest) → helm-validate (lint + template) → "
              "build-backend / build-frontend (Docker to GHCR, PR-safe -- no push) → deploy (gated to "
              "main, helm upgrade + helm test smoke check).", accent=BLUE, body_size=10)
    icon_card(s, Inches(4.14), Inches(3.45), Inches(3.95), Inches(2.2), "Webhook Intake",
              "GitHub push events verified via HMAC (X-Hub-Signature-256) automatically enqueue a scan -- "
              "no manual trigger required for a connected repository.", accent=RED, body_size=10)
    icon_card(s, Inches(8.28), Inches(3.45), Inches(3.95), Inches(2.2), "Feeds Existing Tooling",
              "SARIF export is directly consumable by GitHub/GitLab code-scanning tabs; Slack, email, and "
              "webhook notifications fire on scan completion or failure.", accent=GOLD, body_size=10)

    source_note(s, Inches(LEFT), Inches(5.85), Inches(CW),
                "Today, KAVACH triggers a scan automatically on push and exports SARIF for existing "
                "code-scanning UIs; a native branch-protection “block the merge” gate is not yet a "
                "shipped feature -- see Future Scope.")

    set_notes(s, "Two things are true here and we want to keep them separate. First, KAVACH has its own "
                 "real CI/CD pipeline -- GitHub Actions runs pytest, Helm lint and template validation, "
                 "then builds and pushes Docker images to GHCR only outside pull requests, and deploys via "
                 "Helm only on main, finishing with a helm-test smoke check. Second, and separately, KAVACH "
                 "integrates into a customer's CI/CD by accepting a verified GitHub push webhook that "
                 "automatically enqueues a scan, and by exporting SARIF that plugs straight into GitHub or "
                 "GitLab's existing code-scanning tab, plus Slack/email/webhook notifications on completion. "
                 "We're not claiming a merge-blocking gate exists today -- it doesn't yet -- so we've called "
                 "that out honestly rather than overstating it.")


# ------------------------------------------------------------------ 19 ----
def slide_19_demo(prs):
    s = add_slide(prs, "Live Demonstration Flow", 19, TOTAL)
    kicker_title(s, "Live Demonstration Flow", "What We'll Show, Step by Step")

    steps = [
        ("1", "Repository Upload", "Submit a sandbox payload (low/medium/high risk) or a real repo URL"),
        ("2", "Scan", "Watch all 9 scanners run in parallel with live WebSocket progress"),
        ("3", "Risk Score", "Findings roll up into a Banking Risk Score with risk band"),
        ("4", "AI", "Open a finding -- automated explanation + AI Assistant grounded Q&A"),
        ("5", "Compliance", "Live RBI / PCI DSS / SWIFT CSP PASS/FAIL, with evidence"),
        ("6", "Executive PDF", "Generate the board-facing report on demand"),
        ("7", "Dashboard", "Risk trend, severity distribution, compliance posture"),
        ("8", "Mobile App", "The same scan's status, checked from the Flutter app"),
    ]
    w = (CW - 0.14 * 3) / 4
    for i, (n, t, d) in enumerate(steps):
        x = LEFT + (i % 4) * (w + 0.14)
        y = 2.05 + (i // 4) * 1.85
        numbered_node(s, Inches(x), Inches(y), Inches(w), Inches(1.7), n, t, d, accent=GOLD, body_size=9.5)
        if i % 4 < 3:
            connector(s, Inches(x + w + 0.02), Inches(y + 0.85), Inches(x + w + 0.12), Inches(y + 0.85), color=GOLD, weight=1.5)

    set_notes(s, "This is the exact walkthrough we intend to run live: upload a repository (or one of the "
                 "bundled sandbox payloads at low/medium/high risk), watch all nine scanners fan out with "
                 "real WebSocket progress, land on a Banking Risk Score with its risk band, open a finding "
                 "to see both the automatic explanation and the grounded AI Assistant, check live RBI/PCI/"
                 "SWIFT compliance status with evidence, generate an executive PDF on demand, look at the "
                 "dashboard's risk trend and compliance posture, and finally check that same scan's status "
                 "from the Flutter mobile app to prove it's one backend, not two products.")


# ------------------------------------------------------------------ 20 ----
def slide_20_results(prs):
    s = add_slide(prs, "Results", 20, TOTAL)
    kicker_title(s, "Results", "What The Platform Demonstrably Does")

    stats = [
        ("9 → 1", "Nine scanners' raw findings deduplicated into one correlated, enriched "
                        "UnifiedFinding set per scan -- CVE+package, then file+line+category matching.", BLUE),
        ("~100-200ms", "RAG p50 end-to-end retrieval latency with no LLM call -- embedding ~9ms, vector "
                         "search ~11ms, rerank 15-55ms (measured, see docs/production_hardening.md).", PURPLE_LT),
        ("7 formats", "Executive PDF, Technical PDF, SARIF, CycloneDX SBOM, unified JSON, compliance JSON, "
                        "CSV -- generated automatically per completed scan.", GOLD),
        ("46.6 req/s", "Sustained RAG throughput under a 100-request/3-account load test -- p50 175ms, "
                         "p95 640ms (measured, see docs/production_hardening.md).", GREEN),
    ]
    w = (CW - 0.18 * 3) / 4
    for i, (val, label, color) in enumerate(stats):
        x = LEFT + i * (w + 0.18)
        stat_tile(s, Inches(x), Inches(2.1), Inches(w), Inches(2.75), val, label, accent=color)

    bullets(s, Inches(LEFT), Inches(5.15), Inches(CW), Inches(1.65), [
        ("Better prioritization: ", "a 7-factor Banking Risk Score replaces flat CVSS triage."),
        ("Compliance visibility: ", "live PASS/FAIL against RBI/PCI DSS/SWIFT CSP with evidence, not a "
                                       "quarterly spreadsheet exercise."),
        ("Developer productivity: ", "one scan produces every report format a security, compliance, or "
                                        "engineering stakeholder needs -- no separate tools to reconcile."),
    ], size=12, gap=0.1)

    source_note(s, Inches(LEFT), Inches(6.85), Inches(CW),
                "We report only measured figures from the repository's own benchmark/audit docs -- no "
                "invented percentage improvements.")

    set_notes(s, "We're deliberately only citing numbers we can point at in the repository. The aggregation "
                 "layer collapses nine scanners' overlapping raw findings into one deduplicated set per scan. "
                 "The RAG layer's own production-hardening benchmark shows p50 end-to-end retrieval around "
                 "100-200 milliseconds with no LLM call, and sustains roughly 46.6 requests per second under "
                 "load with p50 175ms. Every completed scan produces all seven report formats automatically. "
                 "Beyond the hard numbers: prioritization improves because BRS replaces flat CVSS, compliance "
                 "visibility improves because PASS/FAIL is live and evidenced rather than a quarterly "
                 "spreadsheet, and developer productivity improves because one scan produces every format a "
                 "stakeholder needs. We're not going to put a fabricated '60% faster' number on this slide -- "
                 "we don't have a controlled before/after study, and the jury will respect the honesty more "
                 "than the number.")


# ------------------------------------------------------------------ 21 ----
def slide_21_future(prs):
    s = add_slide(prs, "Future Scope", 21, TOTAL)
    kicker_title(s, "Future Scope", "Realistic Next Steps -- Grounded In What's Already Built")

    textbox(s, Inches(LEFT), Inches(2.02), Inches(CW), Inches(0.3), "IN ACTIVE DEVELOPMENT",
            size=11.5, color=GOLD, bold=True, letter_spacing=1.0)
    bullets(s, Inches(LEFT), Inches(2.35), Inches(CW), Inches(1.15), [
        "Wiring the Flutter app's live scan-progress screen to the existing WebSocket endpoint (currently polling)",
        "A cross-scan findings/compliance rollup endpoint, so both clients can show portfolio-wide views",
        "Full SAML 2.0 assertion validation (currently a scaffolded 503, pending an XML-security toolkit)",
    ], size=11.5, marker_color=GOLD, gap=0.08)

    textbox(s, Inches(LEFT), Inches(3.75), Inches(5.9), Inches(0.3), "ON THE ROADMAP",
            size=11.5, color=BLUE_LT, bold=True, letter_spacing=1.0)
    bullets(s, Inches(LEFT), Inches(4.08), Inches(5.9), Inches(2.7), [
        "Notifications inbox API + mobile push registration (delivery is outbound-only today)",
        "Knowledge Graph linking findings, CWEs, and compliance clauses",
        "Root Cause Intelligence -- trace a finding class back to the introducing commit",
        "Security Copilot -- propose (not silently apply) a remediation diff",
        "Enterprise integrations -- GitLab/Bitbucket inbound webhooks, PR-scoped scans, Jira/ServiceNow, SIEM export",
        "Expanded compliance coverage -- ISO 27001, SOC 2",
    ], size=11, marker_color=BLUE_LT, gap=0.08)

    textbox(s, Inches(7.0), Inches(3.75), Inches(5.78), Inches(0.3), "PROPOSED DIRECTION (VISION)",
            size=11.5, color=SLATE, bold=True, letter_spacing=1.0)
    bullets(s, Inches(7.0), Inches(4.08), Inches(5.78), Inches(2.7), [
        "IDE plugin -- surface findings at the point of writing code, not just at scan time",
        "Continuous-learning RAG -- incorporate analyst feedback signals already being collected today",
        "Multi-bank / multi-tenant SaaS packaging for smaller banks and NBFCs",
    ], size=11, marker_color=SLATE, gap=0.1)
    source_note(s, Inches(7.0), Inches(5.3), Inches(5.78),
                "Labeled as vision, not in-progress -- not yet started in the codebase.")

    set_notes(s, "We split this into three honesty tiers. In active development, straight from the README's "
                 "own roadmap: wiring the mobile app's live scan progress to the WebSocket endpoint it "
                 "doesn't yet use, a cross-scan rollup endpoint, and full SAML assertion validation. On the "
                 "roadmap, still from the same source: a notifications inbox, a knowledge graph linking "
                 "findings to CWEs and compliance clauses, root-cause intelligence, a security copilot that "
                 "proposes rather than silently applies a fix, deeper enterprise integrations, and expanded "
                 "compliance coverage. The third column -- IDE plugin, continuous-learning RAG, multi-bank "
                 "SaaS -- is explicitly labeled as vision, not in-progress work, because none of it exists in "
                 "the codebase today. We'd rather under-promise on this slide than have a juror open the "
                 "repo and not find it.")


# ------------------------------------------------------------------ 22 ----
def slide_22_thanks(prs):
    s = add_slide(prs, "", 22, TOTAL)
    shield_mark(s, Inches(5.67), Inches(1.15), Inches(1.0), fill=GOLD)
    textbox(s, Inches(3.0), Inches(2.35), Inches(7.33), Inches(0.8), "Thank You",
            size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, Inches(3.0), Inches(3.2), Inches(7.33), Inches(0.5), "KAVACH -- AI-Powered DevSecOps Platform for Banking",
            size=15, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

    card(s, Inches(3.9), Inches(4.05), Inches(5.53), Inches(1.55), fill=PANEL, line=BORDER)
    textbox(s, Inches(4.2), Inches(4.22), Inches(5), Inches(0.3), "TEAM KAVACH", size=11, color=GOLD, bold=True, align=PP_ALIGN.CENTER, letter_spacing=1.2)
    textbox(s, Inches(4.2), Inches(4.55), Inches(5), Inches(0.4), " · ".join(MEMBERS), size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, Inches(4.2), Inches(4.92), Inches(5), Inches(0.35), INSTITUTE, size=11.5, color=SLATE, align=PP_ALIGN.CENTER)
    textbox(s, Inches(4.2), Inches(5.22), Inches(5), Inches(0.3), GITHUB, size=11, color=BLUE_LT, align=PP_ALIGN.CENTER)

    qr = rect(s, Inches(10.2), Inches(4.05), Inches(1.55), Inches(1.55), fill=PANEL, line=BORDER_LT, radius=0.06, dashed=True)
    tf = qr.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "QR\n(GitHub repo)"
    r.font.size = Pt(10.5); r.font.color.rgb = SLATE_DIM; r.font.name = FONT

    textbox(s, Inches(3.0), Inches(6.0), Inches(7.33), Inches(0.5), "Questions?", size=20, color=GOLD_LT,
            bold=True, align=PP_ALIGN.CENTER)

    set_notes(s, "Thank you. We're Team Kavach -- Harshith B, Tejashwan Gangishetty, and Pratham Lal, from "
                 "Vellore Institute of Technology. Everything in this deck is grounded in the actual "
                 "repository -- github.com/TROJAN1HAMMER/KAVACH -- and every claim we made about what's "
                 "implemented versus what's a placeholder or future work is documented in the fact-checklist "
                 "we're submitting alongside this deck. Happy to take questions, and happy to open any file "
                 "in the codebase live if the jury wants to verify a specific claim.")


def build():
    prs = new_deck()
    for fn in [slide_01_title, slide_02_about, slide_03_problem, slide_04_industry, slide_05_solution,
               slide_06_comparison, slide_07_architecture, slide_08_pipeline, slide_09_scanners,
               slide_10_ase, slide_11_brs, slide_12_ai, slide_13_rag, slide_14_compliance,
               slide_15_rbac, slide_16_mobile, slide_17_frontend, slide_18_cicd, slide_19_demo,
               slide_20_results, slide_21_future, slide_22_thanks]:
        fn(prs)
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "KAVACH_PSB_Hackathon_2026.pptx")
    prs.save(out_path)
    print("Saved:", out_path)


if __name__ == "__main__":
    build()
